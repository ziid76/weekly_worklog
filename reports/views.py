from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from django.contrib.auth.mixins import LoginRequiredMixin
from django.views.generic import ListView, DetailView
from django.contrib import messages
from django.db.models import Q
from django.core.paginator import Paginator
from django.http import JsonResponse, HttpResponse
from django.conf import settings
from collections import defaultdict
from copy import deepcopy
import datetime
import re
from openpyxl import Workbook
from pptx import Presentation
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from pptx.util import Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.dml.color import RGBColor
from .models import WeeklyReport, WeeklyReportComment, WeeklyReportPersonalComment, ReportReview, TeamPerformanceAnalysis
from mailing.text_formatter import format_review_content
from worklog.models import Worklog
from teams.models import Team, TeamMembership
from accounts.models import UserProfile
from django.contrib.auth.models import User
from .forms import WeeklyReportCommentForm, WeeklyReportPersonalCommentForm
import html
from django.utils.html import strip_tags
from django.template.loader import render_to_string
from copy import deepcopy
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

class WeeklyReportListView(LoginRequiredMixin, ListView, ):
    model = WeeklyReport
    template_name = 'reports/weekly_report_list.html'
    context_object_name = 'reports'
    paginate_by = 10

    def get_queryset(self):
        return WeeklyReport.objects.filter(team=self.request.user.profile.primary_team.id).order_by('-year', '-week_number')

@login_required
def weekly_report_detail(request, id):
    """주간 리포트 상세 보기"""
    report = get_object_or_404(WeeklyReport, id=id)
    
    
    # 팀이 지정된 경우 해당 팀의 워크로그만, 아니면 전체
    if report.team:
        team = get_object_or_404(Team, id=report.team.id)
        # 팀 멤버들의 워크로그만 가져오기
        team_members = team.members.all()
        worklogs = Worklog.objects.filter(
            year=report.year, 
            week_number=report.week_number,
            author__in=team_members
        ).select_related('author', 'author__profile').order_by('display_order', 'author__profile__last_name_ko', 'author__username')
        

    else:
        # 전체 워크로그
        worklogs = Worklog.objects.filter(
            year=report.year, 
            week_number=report.week_number
        ).select_related('author', 'author__profile').order_by('display_order', 'author__profile__last_name_ko', 'author__username')

    
    # 댓글 처리
    if request.method == 'POST':
        form = WeeklyReportCommentForm(request.POST)
        if form.is_valid():
            comment = form.save(commit=False)
            comment.report = report
            comment.author = request.user
            comment.save()
            messages.success(request, '코멘트가 추가되었습니다.')
            return redirect('weekly_report_detail', id=report.id)
    else:
        form = WeeklyReportCommentForm()

    personal_comment_form = WeeklyReportPersonalCommentForm()

    # 작성자별 워크로그 및 업무 코멘트 그룹화
    personal_comments = report.personal_comments.select_related(
        'target_user__profile',
        'created_by__profile'
    ).order_by('created_at')

    comments_by_user = defaultdict(list)
    for comment in personal_comments:
        comments_by_user[comment.target_user_id].append(comment)

    # 해당 주차의 AI 리뷰가 존재하는 사용자 ID 조회
    review_user_ids = set(ReportReview.objects.filter(
        year=report.year,
        week_number=report.week_number
    ).values_list('user', flat=True))

    worklog_entries = []
    for worklog in worklogs:
        author_profile = getattr(worklog.author, 'profile', None)
        author_name = author_profile.get_korean_name if author_profile else worklog.author.username

        worklog_entries.append({
            'author': worklog.author,
            'author_id': worklog.author_id,
            'author_name': author_name,
            'worklog': worklog,
            'personal_comments': comments_by_user.get(worklog.author_id, []),
            'has_ai_review': worklog.author_id in review_user_ids,
        })

    user_role = None
    if hasattr(request.user, 'profile'):
        user_role = request.user.profile.current_team_role

    can_manage_personal_comments = request.user.is_staff or request.user.is_superuser or user_role in ('leader', 'admin')

    context = {
        'report': report,
        'worklogs': worklogs,
        'worklog_entries': worklog_entries,
        'worklog_count': len(worklog_entries),
        'form': form,
        'comments': report.comments.select_related('author__profile').all(),
        'personal_comment_form': personal_comment_form,
        'can_manage_personal_comments': can_manage_personal_comments,
        'year': report.year,
        'week_number': report.week_number,
        'selected_team': report.team,
    }
    return render(request, 'reports/weekly_report_detail.html', context)


@login_required
def add_personal_comment(request, report_id):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': '잘못된 요청입니다.'}, status=405)

    report = get_object_or_404(WeeklyReport, id=report_id)

    user_role = None
    if hasattr(request.user, 'profile'):
        user_role = request.user.profile.current_team_role

    if not (request.user.is_staff or request.user.is_superuser or user_role in ('leader', 'admin')):
        return JsonResponse({'success': False, 'error': '업무 코멘트를 작성할 권한이 없습니다.'}, status=403)

    target_user_id = request.POST.get('target_user_id')
    if not target_user_id:
        return JsonResponse({'success': False, 'error': '대상 사용자가 지정되지 않았습니다.'}, status=400)

    target_user = get_object_or_404(User, id=target_user_id)

    if not Worklog.objects.filter(
        year=report.year,
        week_number=report.week_number,
        author=target_user
    ).exists():
        return JsonResponse({'success': False, 'error': '해당 사용자의 주간업무가 존재하지 않습니다.'}, status=400)

    form = WeeklyReportPersonalCommentForm(request.POST)
    if not form.is_valid():
        return JsonResponse({'success': False, 'error': '코멘트 내용을 입력해주세요.'}, status=400)

    personal_comment = form.save(commit=False)
    personal_comment.report = report
    personal_comment.target_user = target_user
    personal_comment.created_by = request.user
    personal_comment.save()

    updated_comments = report.personal_comments.filter(
        target_user=target_user
    ).select_related('created_by__profile').order_by('created_at')

    comments_html = render_to_string(
        'reports/personal_comment_list.html',
        {'comments': updated_comments},
        request=request
    )

    return JsonResponse({
        'success': True,
        'html': comments_html,
        'target_user_id': target_user_id,
    })

@login_required

@login_required
def team_worklog_summary(request):
    """팀별 워크로그 요약"""
    year = request.GET.get('year')
    week_number = request.GET.get('week_number')
    team_id = request.GET.get('team')
    
    if not year or not week_number:
        today = datetime.date.today()
        year, week_number, _ = today.isocalendar()
    else:
        year = int(year)
        week_number = int(week_number)
    
    # 팀별로 워크로그 필터링
    if team_id:
        team = get_object_or_404(Team, id=team_id)
        team_members = team.members.all()
        worklogs = Worklog.objects.filter(
            year=year, 
            week_number=week_number,
            author__in=team_members
        ).select_related('author')
        selected_team = team
    else:
        worklogs = Worklog.objects.filter(
            year=year, 
            week_number=week_number
        ).select_related('author')
        selected_team = None
    
    # 팀별로 그룹화 (실제 팀 멤버십 기준)
    team_summary = defaultdict(list)
    
    for worklog in worklogs:
        try:
            author_name = worklog.author.profile.get_korean_name
        except:
            author_name = worklog.author.username
        
        # 사용자의 팀 정보 가져오기
        user_teams = TeamMembership.objects.filter(user=worklog.author).select_related('team')
        
        if user_teams.exists():
            for membership in user_teams:
                team_summary[membership.team.name].append({
                    'author': author_name,
                    'worklog': worklog,
                    'role': membership.get_role_display()
                })
        else:
            team_summary['미분류'].append({
                'author': author_name,
                'worklog': worklog,
                'role': '일반'
            })
    
    # 연도 및 주차 선택 옵션
    current_year = datetime.date.today().year
    year_choices = list(range(current_year - 2, current_year + 3))
    week_choices = list(range(1, 54))
    
    # 팀 목록
    teams = Team.objects.all()
    
    context = {
        'year': year,
        'week_number': week_number,
        'team_summary': dict(team_summary),
        'week_start': datetime.date.fromisocalendar(year, week_number, 1),
        'week_end': datetime.date.fromisocalendar(year, week_number, 1) + datetime.timedelta(days=6),
        'year_choices': year_choices,
        'week_choices': week_choices,
        'teams': teams,
        'selected_team': selected_team,
    }
    
    return render(request, 'reports/team_worklog_summary.html', context)

@login_required
def worklog_summary_popup(request, year, week_number):
    """작성자별 워크로그 집계 팝업"""
    team_id = request.GET.get('team')
    
    # 팀이 지정된 경우 해당 팀의 워크로그만, 아니면 전체
    if team_id:
        team = get_object_or_404(Team, id=team_id)
        # 팀 멤버들의 워크로그만 가져오기
        team_members = team.members.all()
        worklogs = Worklog.objects.filter(
            year=year, 
            week_number=week_number,
            author__in=team_members
        ).select_related('author', 'author__profile').order_by('display_order', 'author__profile__last_name_ko', 'author__username')
        
        report = get_object_or_404(WeeklyReport, year=year, week_number=week_number, team=team)
    else:
        # 전체 워크로그
        worklogs = Worklog.objects.filter(
            year=year, 
            week_number=week_number
        ).select_related('author', 'author__profile').order_by('display_order', 'author__profile__last_name_ko', 'author__username')
        report = get_object_or_404(WeeklyReport, year=year, week_number=week_number, team=None)
    
    # 작성자별 워크로그 그룹화
    worklog_by_author = {}
    for worklog in worklogs:
        # UserProfile이 있으면 한국식 이름 사용, 없으면 username 사용
        try:
            author_name = worklog.author.profile.get_korean_name
        except:
            author_name = worklog.author.username
        
        worklog_by_author[author_name] = worklog
    
    context = {
        'report': report,
        'worklog_by_author': worklog_by_author,
    }
    
    return render(request, 'reports/worklog_summary_popup.html', context)

@login_required
def confirm_closing_api(request):
    try:
        report = get_object_or_404(WeeklyReport, pk=request.POST.get('pk'))
        report.editable = False
        report.save()

        return JsonResponse({
            'success': True,
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        })

@login_required
def export_weekly_report_excel(request, id):
    """주간 리포트 Excel 내보내기"""
    report = get_object_or_404(WeeklyReport, id=id)
    
    # 워크로그 데이터 가져오기
    worklogs = Worklog.objects.filter(
        year=report.year,
        week_number=report.week_number,
        author__teams=report.team
    ).select_related('author', 'author__profile').order_by('display_order', 'author__profile__last_name_ko', 'author__username')
    
    worklog_by_author = {}
    for worklog in worklogs:
        author_name = worklog.author.profile.display_name if worklog.author.profile else worklog.author.username
        worklog_by_author[author_name] = worklog
    
    # Excel 워크북 생성
    wb = Workbook()
    ws = wb.active
    ws.title = f"{report.year}년 {report.week_number}주차 주간보고서"
    
    # 스타일 정의
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
    border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    # 헤더 작성
    ws['A1'] = "담당자"
    ws['B1'] = f"금주 실적 ({report.week_start_date.strftime('%m월 %d일')} ~ {report.week_end_date.strftime('%m월 %d일')})"
    ws['C1'] = f"차주 계획 ({report.next_week_start_date.strftime('%m월 %d일')} ~ {report.next_week_end_date.strftime('%m월 %d일')})"
    
    # 헤더 스타일 적용
    for col in ['A1', 'B1', 'C1']:
        ws[col].font = header_font
        ws[col].fill = header_fill
        ws[col].alignment = Alignment(horizontal='center', vertical='center')
        ws[col].border = border
    
    # 컬럼 너비 설정
    ws.column_dimensions['A'].width = 15
    ws.column_dimensions['B'].width = 50
    ws.column_dimensions['C'].width = 50
    
    # 데이터 작성
    row = 2
    for author_name, worklog in worklog_by_author.items():
        ws[f'A{row}'] = author_name
        
        # 마크다운 텍스트를 일반 텍스트로 변환
        this_week_text = clean_markdown_text(worklog.this_week_work) if worklog.this_week_work else "작성된 내용이 없습니다."
        next_week_text = clean_markdown_text(worklog.next_week_plan) if worklog.next_week_plan else "작성된 계획이 없습니다."
        
        ws[f'B{row}'] = this_week_text
        ws[f'C{row}'] = next_week_text
        
        # 셀 스타일 적용
        for col in ['A', 'B', 'C']:
            cell = ws[f'{col}{row}']
            cell.border = border
            cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
        
        row += 1
    
    # Excel 파일을 메모리 버퍼에 저장
    import io
    from django.http import FileResponse
    
    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    
    # 팀명과 월/주차 정보로 파일명 생성
    team_name = report.team.name if report.team else "전체"
    month = report.week_start_date.month
    week_in_month = ((report.week_start_date.day - 1) // 7) + 1
    filename = f"{team_name}_주간보고서_{month}월_{week_in_month}주차.xlsx"
    
    response = FileResponse(
        buffer,
        as_attachment=True,
        filename=filename,
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    return response






    return response


@login_required
def get_ai_review_content(request, report_id, user_id):
    """특정 주간보고의 사용자 AI 리뷰 내용을 가져옵니다 (AJAX)."""
    if not request.headers.get('x-requested-with') == 'XMLHttpRequest':
        return JsonResponse({'success': False, 'message': '잘못된 요청입니다.'}, status=400)

    report = get_object_or_404(WeeklyReport, id=report_id)
    target_user = get_object_or_404(User, id=user_id)
    
    # 해당 연도/주차 및 사용자의 가장 최근 리뷰 조회
    review = ReportReview.objects.filter(
        year=report.year,
        week_number=report.week_number,
        user=target_user
    ).order_by('-created_at').first()
    
    if not review:
        return JsonResponse({'success': False, 'message': '해당 사용자의 AI 리뷰가 아직 존재하지 않습니다.'})
    
    # 이메일 템플릿에 사용할 데이터 포맷팅
    formatted_review = format_review_content(review.review_content)
    
    context = {
        'review': formatted_review,
        'month_week_display': review.month_week_display,
        'user': target_user,
        'year': report.year,
        'week_number': report.week_number,
        'site_url': settings.SITE_URL,
    }
    
    # 이메일 템플릿을 사용하여 HTML 생성 (inline-style 포함)
    html_content = render_to_string('emails/review_notification.html', context)
    
    return JsonResponse({'success': True, 'html': html_content})


@login_required
def export_weekly_report_pptx(request, id):
    """주간 리포트를 PowerPoint 파일로 내보냅니다."""
    report = get_object_or_404(WeeklyReport, id=id)

    # 팀이 지정된 경우 해당 팀의 워크로그만, 아니면 전체
    if report.team:
        team = get_object_or_404(Team, id=report.team.id)
        # 팀 멤버들의 워크로그만 가져오기
        team_members = team.members.all()
        worklogs = Worklog.objects.filter(
            year=report.year, 
            week_number=report.week_number,
            author__in=team_members
        ).select_related('author', 'author__profile').order_by('display_order', 'author__profile__last_name_ko', 'author__username')
    else:
        # 전체 워크로그
        worklogs = Worklog.objects.filter(
            year=report.year, 
            week_number=report.week_number
        ).select_related('author', 'author__profile').order_by('display_order', 'author__profile__last_name_ko', 'author__username')

    entries = []
    for worklog in worklogs:
        profile = getattr(worklog.author, 'profile', None)
        author_name = profile.get_korean_name if profile else worklog.author.username
        meta_parts = []
        if profile and getattr(profile, 'position', None):
            meta_parts.append(profile.position)
        if profile and getattr(profile, 'department_display', None):
            meta_parts.append(profile.department_display)
        meta_text = " / ".join(meta_parts)

        this_week = clean_markdown_text_for_pptx(worklog.this_week_work) if worklog.this_week_work else "작성된 내용이 없습니다."
        next_week = clean_markdown_text_for_pptx(worklog.next_week_plan) if worklog.next_week_plan else "작성된 계획이 없습니다."

        entries.append(
            {
                "author": author_name,
                "this_week": this_week,
                "next_week": next_week,
            }
        )

    if not entries:
        entries.append(
            {
                "author": "작성자 없음",
                "this_week": "작성된 내용이 없습니다.",
                "next_week": "작성된 계획이 없습니다.",
            }
        )

    month_week_display = getattr(report, "month_week_display", None)
    if not month_week_display:
        week_in_month = ((report.week_start_date.day - 1) // 7) + 1
        month_week_display = f"{report.week_start_date.month}월 {week_in_month}주차"
    team_name = report.team.name if report.team else "전체"

    template_path = settings.BASE_DIR / "reports/files/templates.pptfile"
    prs = Presentation(str(template_path))
    base_slide = prs.slides[0]

    def get_table(target_slide):
        for shape in target_slide.shapes:
            if shape.has_table:
                return shape.table
        return None

    base_table = get_table(base_slide)
    if base_table is None:
        return HttpResponse("템플릿에 표가 없습니다.", status=500)

    if len(base_table.rows) < 2:
        return HttpResponse("템플릿에 데이터 행이 필요합니다.", status=500)

    template_row_xml = deepcopy(base_table._tbl.tr_lst[1])
    capacity = max(len(base_table.rows) - 1, 1)

    template_shapes_xml = [deepcopy(shape._element) for shape in base_slide.shapes]
    template_relationships = []
    for r_id, rel in base_slide.part.rels.items():  # (rId, rel) 튜플
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue
        # 발표자 노트는 복사하지 않는 편이 안전
        if "notesSlide" in str(rel.reltype):
            continue
        if getattr(rel, "is_external", False):
            template_relationships.append((rel.reltype, rel.target_ref, True))   # URL 문자열
        else:
            template_relationships.append((rel.reltype, rel.target_part, False)) # 내부 파트

    next_week_start = getattr(report, "next_week_start_date", None)
    next_week_end = getattr(report, "next_week_end_date", None)
    if not next_week_start or not next_week_end:
        next_week_start = report.week_end_date + datetime.timedelta(days=3)
        next_week_end = next_week_start + datetime.timedelta(days=4)

    header_labels = [
        "담당",
        f"금주 실적 ({report.week_start_date:%m월 %d일} ~ {report.week_end_date:%m월 %d일})",
        f"차주 계획 ({next_week_start:%m월 %d일} ~ {next_week_end:%m월 %d일})",
    ]

    def clone_template_slide():
        new_slide = prs.slides.add_slide(base_slide.slide_layout)

        # 기본 placeholder 제거
        for shape in list(new_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)

        # 레이아웃 외 관계 제거
        for r_id, rel in list(new_slide.part.rels.items()):
            if rel.reltype != RT.SLIDE_LAYOUT:
                del new_slide.part.rels[r_id]

        # 도형 XML 복사
        for shape_xml in template_shapes_xml:
            new_slide.shapes._spTree.insert_element_before(deepcopy(shape_xml), "p:extLst")

        # 관계 복사 (_add_relationship 사용, rId는 라이브러리에 맡김)
        for reltype, target, is_external in template_relationships:
            if hasattr(new_slide.part.rels, "_add_relationship"):
                new_slide.part.rels._add_relationship(reltype, target, is_external)
            else:
                # 0.6.19 이하 하위호환 (rId 없이 호출)
                new_slide.part.rels.add_relationship(reltype, target)

        return new_slide
        
    def fit_cell_text(cell, text, *, max_size=11, min_size=9, bold=False, font_color=None):
        """표 셀 텍스트를 자동 축소하여 맞춤 - bold 마크다운 처리"""
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = tf.paragraphs[0]
        p.line_spacing = 1.2  # 줄간격 1.2배수로 설정
        p.space_after = Pt(0)
        
        if not text:
            text = ""
        
        # **bold** 마크다운 처리
        parts = re.split(r'(\*\*.*?\*\*)', text)
        
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                # bold 텍스트
                bold_text = part[2:-2]  # ** 제거
                run = p.add_run()
                run.text = bold_text
                run.font.bold = True
                run.font.size = Pt(max_size)
                if font_color:
                    run.font.color.rgb = font_color
            elif part:
                # 일반 텍스트
                run = p.add_run()
                run.text = part
                run.font.bold = bold
                run.font.size = Pt(max_size)
                if font_color:
                    run.font.color.rgb = font_color
        
        # 줄 수 계산
        lines = max((text or "").count("\n") + 1, 1)
        return lines

    def tighten_cell_layout(cell):
        """표 셀 여백과 줄간격을 최소화"""
        # 표 셀 여백 최소화
        cell.margin_left = Pt(4)   # 너무 작으면 텍스트가 잘림
        cell.margin_right = Pt(4)
        cell.margin_top = Pt(3)
        cell.margin_bottom = Pt(3)

    def set_cell_text(cell, text, *, bold=False, max_size=11, min_size=9):
        """개선된 셀 텍스트 설정"""
        tighten_cell_layout(cell)
        lines = fit_cell_text(cell, text, max_size=max_size, min_size=min_size, bold=bold)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        return lines

    def optimize_table_layout(table):
        """표 레이아웃 최적화 - 컬럼 너비 재분배 및 가운데 정렬"""
        if not table or len(table.columns) < 3:
            return
            
        # 슬라이드 너비 기준으로 표 너비 설정 (95% 사용하여 좌우 여백 최소화 및 균형 확보)
        slide_width = prs.slide_width
        table_width = int(slide_width * 0.95)
        
        # 최적 비율: 담당(10%) : 금주실적(45%) : 차주계획(45%)
        ratios = [0.1, 0.45, 0.45]
        
        # 각 컬럼에 절대 너비 적용
        for i, ratio in enumerate(ratios):
            if i < len(table.columns):
                table.columns[i].width = int(table_width * ratio)

    def populate_table(target_slide, chunk):
        """개선된 표 채우기 함수"""
        table = get_table(target_slide)
        if table is None:
            return False
            
        # 표 레이아웃 최적화
        optimize_table_layout(table)
        
        # 표를 가운데 정렬
        for shape in target_slide.shapes:
            if shape.has_table:
                slide_width = prs.slide_width
                table_width = sum(col.width for col in shape.table.columns)
                # 가운데 정렬을 위한 left 위치 계산
                shape.left = (slide_width - table_width) // 2
                break
            
        # 기존 행 제거 (헤더 제외)
        while len(table.rows) > 1:
            table._tbl.remove(table._tbl.tr_lst[-1])
            
        # 새 행 추가
        for _ in chunk:
            table._tbl.append(deepcopy(template_row_xml))
            
        # 헤더 설정 (3개 컬럼만 사용)
        header_row = table.rows[0]
        for i, label in enumerate(header_labels):
            if i < len(header_row.cells):
                tighten_cell_layout(header_row.cells[i])
                fit_cell_text(header_row.cells[i], label, max_size=12, min_size=12, bold=True, font_color=RGBColor(255, 255, 255))
        
        # 4번째 컬럼이 있다면 숨기기
        if len(header_row.cells) > 3:
            header_row.cells[3].text = ""
            
        table.rows[0].height = Pt(38)
        
        # 본문 데이터 설정 및 높이 계산
        total_height = table.rows[0].height.pt
        # 슬라이드 실제 높이에 맞게 최대 높이 설정 (약 75% 수준)
        slide_height_pt = prs.slide_height / 12700
        slide_body_max = slide_height_pt * 0.75
        
        for row_idx, entry in enumerate(chunk, start=1):
            if row_idx >= len(table.rows):
                break
                
            row = table.rows[row_idx]
            
            # 각 셀 텍스트 설정 및 줄 수 계산 (3개 컬럼만 사용)
            l1 = set_cell_text(row.cells[0], entry["author"], bold=True, max_size=13, min_size=9)
            l2 = set_cell_text(row.cells[1], entry["this_week"], max_size=10, min_size=10)  # 금주실적을 2번째 컬럼에
            l3 = set_cell_text(row.cells[2], entry["next_week"], max_size=10, min_size=10)  # 차주계획을 3번째 컬럼에
            
            # 4번째 컬럼이 있다면 비우기
            if len(row.cells) > 3:
                row.cells[3].text = ""
            
            # 행 높이 계산 (최대 줄 수 기준)
            max_lines = max(l1, l2, l3)
            row_height_pt = 24 + (max_lines - 1) * 12  # 기본 높이와 줄간격 최적화
            row.height = Pt(row_height_pt)
            total_height += row_height_pt
            
            # 슬라이드 높이 초과 검사 (단, 최소 1개의 행은 보장)
            if total_height > slide_body_max and row_idx > 1:
                # 현재 행 제거하고 분할 필요 표시
                table._tbl.remove(table._tbl.tr_lst[-1])
                return False  # 분할 필요
                
        return True  # 성공적으로 완료

    # 동적 슬라이드 분할 처리
    template_slide = base_slide
    slide_index = 0
    entry_index = 0
    
    while entry_index < len(entries):
        # 현재 슬라이드 준비
        slide = template_slide if slide_index == 0 else clone_template_slide()
        
        # 남은 항목들로 청크 생성 (최대 capacity개)
        remaining_entries = entries[entry_index:]
        chunk = remaining_entries[:capacity]
        
        # 표 채우기 시도
        success = populate_table(slide, chunk)
        
        if success:
            # 성공: 전체 청크 처리 완료
            entry_index += len(chunk)
        else:
            # 실패: 청크 크기를 줄여서 재시도
            if len(chunk) > 1:
                # 절반으로 줄여서 재시도
                chunk = chunk[:len(chunk)//2]
                success = populate_table(slide, chunk)
                if success:
                    entry_index += len(chunk)
                else:
                    # 1개씩이라도 처리
                    chunk = chunk[:1]
                    populate_table(slide, chunk)
                    entry_index += 1
            else:
                # 1개 항목도 들어가지 않는 경우 (매우 긴 텍스트)
                entry_index += 1
        
        # 슬라이드 제목 설정
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = f"{month_week_display} {team_name} 주간보고"
        
        # 빈 텍스트 박스 처리
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                existing = shape.text.replace(" ", "").strip()
                if not existing:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0]
                    paragraph.font.size = Pt(12)
                    break
        
        slide_index += 1

    # PPTX 파일을 메모리 버퍼에 저장
    import io
    from django.http import FileResponse
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    filename = f"{team_name}_주간보고_{month_week_display}.pptx"
    
    response = FileResponse(
        buffer,
        as_attachment=True,
        filename=filename,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    return response


def clean_markdown_text_for_pptx(text):
    """PPTX용 텍스트 정리 - 원본 형태 최대한 유지, bold 처리 포함"""
    if not text:
        return ""

    # HTML 태그 처리 (bold는 유지)
    text = text.replace("<br>", "\n").replace("<br/>", "\n").replace("<br />", "\n")
    text = re.sub(r'</p\s*>', '\n', text, flags=re.IGNORECASE)
    
    # strong/b 태그를 **로 변환하여 유지
    text = re.sub(r'<strong[^>]*>(.*?)</strong>', r'**\1**', text, flags=re.IGNORECASE | re.DOTALL)
    text = re.sub(r'<b[^>]*>(.*?)</b>', r'**\1**', text, flags=re.IGNORECASE | re.DOTALL)
    
    # font-weight 스타일 처리 (bolder, bold, 700 등)
    text = re.sub(r'<span[^>]*style="[^"]*font-weight:\s*(?:bolder|bold|700)[^"]*"[^>]*>(.*?)</span>', r'**\1**', text, flags=re.IGNORECASE | re.DOTALL)
    text = re.sub(r'<span[^>]*style=\'[^\']*font-weight:\s*(?:bolder|bold|700)[^\']*\'[^>]*>(.*?)</span>', r'**\1**', text, flags=re.IGNORECASE | re.DOTALL)
    
    # 나머지 HTML 태그 제거
    text = strip_tags(text)
    text = html.unescape(text)
    
    # 기본적인 마크다운만 제거 (** bold는 유지)
    text = re.sub(r'#{1,6}\s*', '', text)
    text = re.sub(r'`(.*?)`', r'\1', text)
    
    # 연속된 빈 줄만 정리
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    
    return text.strip()

def clean_markdown_text(text):
    """마크다운 및 HTML 텍스트를 일반 텍스트로 변환하고 가독성 개선"""
    if not text:
        return ""

    text = text.replace("<br>", "\n").replace("<br/>", "\n").replace("<br />", "\n")
    text = re.sub(r'</p\s*>', '\n', text, flags=re.IGNORECASE)

    # HTML 태그 제거
    text = strip_tags(text)
    
    # HTML 엔티티 디코딩
    text = html.unescape(text)
    
    # 마크다운 문법 제거
    text = re.sub(r'#{1,6}\s*', '', text)  # 헤더
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # 굵은 글씨
    text = re.sub(r'\*(.*?)\*', r'\1', text)  # 기울임
    text = re.sub(r'`(.*?)`', r'\1', text)  # 인라인 코드
    
    # 기존 블릿 기호 제거 (• - * + 등)
    text = re.sub(r'^\s*[•\-\*\+]\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*\d+\.\s+', '', text, flags=re.MULTILINE)  # 번호 리스트
    
    # 긴 문장을 글머리표로 분리 (가독성 개선)
    lines = text.split('\n')
    processed_lines = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # 긴 문장(80자 이상)이고 문장 구분자가 있는 경우 분리
        if len(line) > 80 and any(sep in line for sep in ['. ', ', ', '; ']):
            # 문장 구분자로 분리
            parts = re.split(r'([.;,])\s+', line)
            current_part = ""
            
            for i in range(0, len(parts), 2):
                if i + 1 < len(parts):
                    segment = parts[i] + parts[i + 1]
                else:
                    segment = parts[i]
                
                if len(current_part + segment) > 60 and current_part:
                    processed_lines.append(f"• {current_part.strip()}")
                    current_part = segment
                else:
                    current_part += segment
            
            if current_part.strip():
                processed_lines.append(f"• {current_part.strip()}")
        else:
            processed_lines.append(f"• {line}")
    
    text = '\n'.join(processed_lines)
    
    # 연속된 공백과 줄바꿈 정리
    text = re.sub(r'\n\s*\n', '\n', text)  # 빈 줄 제거
    text = re.sub(r' +', ' ', text)  # 연속 공백 제거
    
    return text.strip()

import json

@login_required
def personal_report_history(request):
    """ 개인별 주간업무 이력 조회 """
    # 현재 로그인한 유저의 팀 정보 가져오기
    try:
        primary_team = request.user.profile.primary_team
        team_members = User.objects.filter(teams=primary_team).select_related('profile').order_by('profile__last_name_ko', 'profile__first_name_ko')
    except AttributeError:
        # 팀 정보가 없는 경우 (예: 관리자)
        primary_team = None
        team_members = User.objects.all().select_related('profile').order_by('profile__last_name_ko', 'profile__first_name_ko')

    # 필터링 파라미터
    selected_user_id = request.GET.get('user_id')
    
    # 날짜 기본값 설정 (최근 5주)
    today = datetime.date.today()
    end_year_default, end_week_default, _ = today.isocalendar()
    start_date_default = today - datetime.timedelta(weeks=4)
    start_year_default, start_week_default, _ = start_date_default.isocalendar()

    start_year = request.GET.get('start_year', str(start_year_default))
    start_week = request.GET.get('start_week', str(start_week_default))
    end_year = request.GET.get('end_year', str(end_year_default))
    end_week = request.GET.get('end_week', str(end_week_default))

    # 연도 선택 옵션
    current_year = datetime.date.today().year
    year_choices = list(range(current_year - 3, current_year + 1))

    # 연도별 주차 정보 생성 (for JavaScript)
    week_data = {}
    for year in year_choices:
        week_data[year] = []
        for week_num in range(1, 54):
            try:
                week_start = datetime.date.fromisocalendar(year, week_num, 1)
                # 주차의 마지막 날이 다음 해로 넘어가는 경우 방지
                if week_start.year != year and week_num > 1:
                    continue
                week_end = week_start + datetime.timedelta(days=4)
                display_text = f"{week_num}주차 ({week_start.strftime('%m.%d')}~{week_end.strftime('%m.%d')})"
                week_data[year].append({'week': week_num, 'display': display_text})
            except ValueError:
                # 53주차가 없는 해 처리
                break

    worklogs_data = []
    if selected_user_id:
        # worklog 쿼리
        worklogs = Worklog.objects.filter(
            author_id=selected_user_id,
        ).order_by('-year', '-week_number')

        # 기간 필터링
        worklogs = worklogs.filter(
            Q(year__gt=int(start_year)) |
            Q(year=int(start_year), week_number__gte=int(start_week))
        )
        worklogs = worklogs.filter(
            Q(year__lt=int(end_year)) |
            Q(year=int(end_year), week_number__lte=int(end_week))
        )

        for log in worklogs:
            # 전주 계획을 가져오기
            prev_week_date = log.week_start_date - datetime.timedelta(days=7)
            prev_year, prev_week, _ = prev_week_date.isocalendar()
            
            previous_worklog = Worklog.objects.filter(
                author_id=selected_user_id,
                year=prev_year,
                week_number=prev_week
            ).first()
            
            worklogs_data.append({
                'worklog': log,
                'previous_week_plan': previous_worklog.next_week_plan if previous_worklog else ''
            })

    context = {
        'team_members': team_members,
        'year_choices': year_choices,
        'week_data_json': json.dumps(week_data),
        'selected_user_id': selected_user_id,
        'start_year': start_year,
        'start_week': start_week,
        'end_year': end_year,
        'end_week': end_week,
        'worklogs_data': worklogs_data,
    }
    return render(request, 'reports/personal_report_history.html', context)


@login_required
def team_performance_list(request):
    """팀 성과 분석 리포트 목록"""
    # 사용자가 속한 팀의 리포트만 보여주거나, 관리자 권한에 따라 전체 조회
    if request.user.is_superuser or request.user.is_staff:
        analyses = TeamPerformanceAnalysis.objects.all()
    else:
        # 일반 사용자는 본인 팀 리포트만
        user_teams = request.user.team_memberships.values_list('team', flat=True)
        analyses = TeamPerformanceAnalysis.objects.filter(team__in=user_teams)
    
    return render(request, 'reports/performance_analysis_list.html', {'analyses': analyses})

@login_required
def team_performance_detail(request, pk):
    """팀 성과 분석 리포트 상세"""
    analysis = get_object_or_404(TeamPerformanceAnalysis, pk=pk)
    
    # 권한 체크 (간단히 팀 멤버인지 확인)
    if not (request.user.is_superuser or request.user.is_staff):
         is_member = request.user.team_memberships.filter(team=analysis.team).exists()
         if not is_member:
             return render(request, '403.html', status=403)
    
    return render(request, 'reports/performance_analysis_detail.html', {'analysis': analysis})